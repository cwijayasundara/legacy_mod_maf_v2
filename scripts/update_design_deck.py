"""Append / rewrite slides in docs/design_deck.pptx to reflect the
large-file handling design + AWS Lambda comprehension wiring.

Three new slides are added and two existing slides are updated in place.
Existing slide 12 serves as the visual template for every new slide
(two-column content under a pipeline row) so the deck stays consistent.

Run as: ``python scripts/update_design_deck.py``.

The script is idempotent-on-text-edit for the existing slides but
appends new slides each run, so delete the generated slides before
re-running if iterating.
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

DECK = Path(__file__).resolve().parents[1] / "docs" / "design_deck.pptx"

# Slide 12 — Stage 3 LLM1 (comprehension) — is the visual template for
# every new "two-column under a pipeline row" slide we add. Cloning its
# XML preserves fonts / colors / geometry exactly.
TEMPLATE_SLIDE_INDEX = 12


def _set_text(shape, text: str) -> None:
    """Overwrite a shape's text while preserving the first run's formatting.

    python-pptx's ``tf.text = ...`` wipes runs and drops styling. We instead
    replace the first paragraph/run's text and clear the rest, which keeps
    the font and color the template baked in.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    paragraphs = text.split("\n")
    # Reuse the first paragraph's first run for styling.
    p0 = tf.paragraphs[0]
    # Wipe extra paragraphs.
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    # Wipe extra runs in p0, keep the first for style.
    runs = list(p0.runs)
    if runs:
        first_run = runs[0]
        for extra_run in runs[1:]:
            extra_run._r.getparent().remove(extra_run._r)
        first_run.text = paragraphs[0]
    else:
        p0.text = paragraphs[0]
    # Add remaining paragraphs, inheriting formatting from p0.
    for line in paragraphs[1:]:
        new_p = copy.deepcopy(p0._p)
        # Strip existing runs on the clone and replace with one plain run.
        for r in list(new_p.findall(".//" + _q("r"))):
            new_p.remove(r)
        # Build a single run using the style inherited from the clone.
        new_run_p = tf.add_paragraph()
        # Mirror the template paragraph's pPr (alignment, bullet etc.).
        pPr_src = p0._p.find(_q("pPr"))
        if pPr_src is not None and new_run_p._p.find(_q("pPr")) is None:
            new_run_p._p.insert(0, copy.deepcopy(pPr_src))
        if runs:
            r_clone = copy.deepcopy(runs[0]._r)
            # Replace text of the cloned run.
            for t in r_clone.findall(_q("t")):
                t.text = line
            # The add_paragraph() already put a placeholder run — replace it.
            for existing_r in list(new_run_p._p.findall(_q("r"))):
                new_run_p._p.remove(existing_r)
            new_run_p._p.append(r_clone)
        else:
            new_run_p.text = line


def _q(tag: str) -> str:
    return f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"


def clone_slide(prs: Presentation, template_index: int) -> object:
    """Clone slide at ``template_index`` and return the new Slide object.

    Copies every shape's XML element into a fresh blank slide, then
    re-attaches relationships so images / chart refs survive. This is
    the standard python-pptx clone recipe — we can't swap whole spTree
    because python-pptx's ``add_slide`` seeds it with layout-derived
    placeholders we need to keep.
    """
    from pptx.oxml.ns import qn

    src = prs.slides[template_index]
    # "Blank" layout is slide_layouts[6] — has no seeded placeholders
    # to collide with the copied shapes.
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    sp_tree = new_slide.shapes._spTree
    for shape in src.shapes:
        new_el = copy.deepcopy(shape.element)
        sp_tree.insert_element_before(new_el, "p:extLst")
    # Re-attach non-external rels (images, fills, etc.).
    for rel in src.part.rels.values():
        if rel.is_external:
            continue
        new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
    return new_slide


def reorder_slide(prs: Presentation, new_slide, insert_after: int) -> None:
    """Move the just-appended slide to right after ``insert_after`` (0-based)."""
    slides = prs.slides
    sldIdLst = slides._sldIdLst
    entries = list(sldIdLst)
    # Pop the last-added entry (our new slide) and insert at position.
    new_entry = entries[-1]
    sldIdLst.remove(new_entry)
    sldIdLst.insert(insert_after + 1, new_entry)


# --------------------------------------------------------------------------- #
# Slide content definitions — one dict per slide to add, keyed by shape name.
# --------------------------------------------------------------------------- #


LARGE_FILE_SLIDE = {
    "title_bar": "PART 2 — PLATFORM ARCHITECTURE",
    "title": "Large-file handling  ·  40k-LOC programs without context overflow",
    "subtitle": (
        "Four layers turn a 500k-token source file into a sequence of bounded "
        "LLM calls. No single prompt ever sees raw source at full size."
    ),
    # Pipeline row labels (replacing KG → Summariser → CRUD → Seam → Spec)
    "pipe_1": "1 · AST decompose\nfile → program /\nsection / paragraph",
    "pipe_2": "2 · Auto-chunk\noversized nodes at\nfunction boundaries",
    "pipe_3": "3 · Map (parallel)\nLLM call per\nchunk / leaf node",
    "pipe_4": "4 · Reduce\nchunk summaries →\nnode `llm_summary`",
    "pipe_5": "5 · Roll up\nnode summaries →\nprogram spec",
    "left_header": "The four-layer design",
    "left_body": (
        "• Layer 1 — AST decomposition (free)\n"
        "Adapters emit program / section / paragraph /\n"
        "field nodes. A 40k-LOC COBOL program becomes\n"
        "10-20 sections and 200-800 paragraphs.\n"
        "• Layer 2 — Semantic sub-chunking (free)\n"
        "When a single paragraph still busts the token\n"
        "budget, the chunker splits at function / class /\n"
        "section boundaries with 50-line overlap.\n"
        "• Layer 3 — Bottom-up map (LLM₁, parallel)\n"
        "Post-order DFS. Leaves first; each chunk gets a\n"
        "3-5 sentence summary. Content-hash cached.\n"
        "• Layer 4 — Reduce (LLM₁ reduce step)\n"
        "Chunk summaries fold into the paragraph's\n"
        "canonical `llm_summary` via the kind prompt.\n"
        "Program-level prompt never sees raw source."
    ),
    "right_header": "Cost math on a single 40k-LOC file",
    "right_body": (
        "• First run at gpt-5.4-mini rates\n"
        "800 paragraphs × ~1.5k in / 300 out ≈ $8.\n"
        "Sub-chunking 5% oversized paragraphs ≈ $1.\n"
        "Section + program summaries ≈ $0.50.\n"
        "Total ≈ $10 per 40k-LOC file, first pass.\n"
        "• Second run (cache hits only)\n"
        "Re-summarizes only nodes whose content hash\n"
        "changed. Unchanged 40k-LOC file ≈ $0.10.\n"
        "• Sampling for 100k+ LOC estates\n"
        "`LEGACY_MOD_COMPREHENSION_SAMPLE_PROGRAMS=N`\n"
        "limits to top-N programs by fan-in centrality.\n"
        "• What LLM₂ sees afterwards\n"
        "The rendered business spec + KG-derived CRUD\n"
        "inputs/outputs — never the raw 40k lines."
    ),
    "footer": (
        "Part 2  ·  Large-file handling  ·  map-reduce summarization + semantic chunking "
        "+ content-hash cache"
    ),
}

CACHE_KNOBS_SLIDE = {
    "title_bar": "PART 2 — PLATFORM ARCHITECTURE",
    "title": "Comprehension cache + operator knobs",
    "subtitle": (
        "Persistent cache, hard caps, centrality sampling — three controls "
        "that keep comprehension cost predictable at estate scale."
    ),
    "pipe_1": "env var\nMAX_SUMMARIES\nhard call cap",
    "pipe_2": "env var\nSAMPLE_PROGRAMS\ntop-N centrality",
    "pipe_3": "persistent cache\ncomprehension_\ncache.jsonl",
    "pipe_4": "chunk_key\nhash-per-chunk\nedit-granular",
    "pipe_5": "budget tracker\nrolls up into\ncost_report.json",
    "left_header": "Operator controls",
    "left_body": (
        "• LEGACY_MOD_COMPREHENSION_MAX_SUMMARIES\n"
        "Hard cap on fresh LLM calls this pass issues.\n"
        "Unset or 0 → comprehension disabled (spec falls\n"
        "back to deterministic-only mode — today's baseline).\n"
        "• LEGACY_MOD_COMPREHENSION_SAMPLE_PROGRAMS\n"
        "Only summarize the top-N programs by centrality\n"
        "(fan-in + fan-out) plus their descendants.\n"
        "Essential for 100k+ LOC estates.\n"
        "• Two knobs compose\n"
        "SAMPLE=20 + MAX=2000 → bounded first-pass\n"
        "triage over the most-referenced programs.\n"
        "• Auto-detects the chat client\n"
        "Reuses the run's LEGACY_MOD_PROVIDER selection\n"
        "(OpenAI / Anthropic / Azure / Ollama).\n"
        "• Comprehension cost lands in the same report\n"
        "Budget tracker tags calls as `comprehension`;\n"
        "by-role breakdown is visible per run."
    ),
    "right_header": "Cache shape + incrementality",
    "right_body": (
        "• On-disk JSONL cache\n"
        "`<workdir>/comprehension_cache.jsonl`\n"
        "Append-only. Dict index rebuilt in memory on load.\n"
        "Survives process restarts + new runs.\n"
        "• Per-node cache key\n"
        "SHA256 of (raw_text | name | kind) + prompt hash.\n"
        "Editing one paragraph invalidates only that entry.\n"
        "• Per-chunk cache key  (new for large files)\n"
        "`SummaryCache.chunk_key(node_id, chunk, prompt_hash)`.\n"
        "Editing one chunk of a 40k-LOC paragraph\n"
        "re-summarizes only that chunk.\n"
        "• Second-run economics\n"
        "Real-run measurement on the AWS corpus:\n"
        "656 cache entries after run 1, second run hits\n"
        "cache on every comprehension call.\n"
        "• Cartridge-agnostic\n"
        "Same cache shape serves COBOL + AWS cartridges."
    ),
    "footer": (
        "Part 2  ·  Comprehension cache + knobs  ·  content-hash keyed, "
        "cartridge-agnostic, budget-tracked"
    ),
}

AWS_PIPELINE_SLIDE = {
    "title_bar": "PART 3 — SAMPLE CARTRIDGE 1",
    "title": "AWS Lambda  ·  full pipeline on disk + pilot results",
    "subtitle": (
        "Every intermediate artifact the platform produces is persisted for "
        "review: AST, KG, business specs, dependency graphs — not just the "
        "migrated code."
    ),
    "pipe_1": "source → AST\nper-unit IR dumps\nast/*.json",
    "pipe_2": "AST → KG\nfull graph snapshot\nkg.json",
    "pipe_3": "KG → LLM₁\nbottom-up summaries\ncache.jsonl",
    "pipe_4": "KG → spec\nbusiness_specs/\n<unit>.md",
    "pipe_5": "KG → diagrams\ngraphs/*.dot .mmd\n.svg (call + data)",
    "left_header": "What lands under  azure_fn_migrated/<repo>/",
    "left_body": (
        "• kg.json\n"
        "Full graph snapshot — every node + edge, with\n"
        "llm_summary populated post-comprehension.\n"
        "• ast/<unit>.json\n"
        "Per-unit flat IR: imports + SDK call sites with\n"
        "file:line granularity. The `AST output`.\n"
        "• business_specs/<unit>.md  (central)\n"
        "Target-style-aware BRD: Purpose (LLM), Inputs /\n"
        "Outputs / Side-effects (KG-CRUD), Paragraphs\n"
        "(per-function LLM summaries).\n"
        "• units/<unit>/business_spec.md  (colocated)\n"
        "Same BRD shipped alongside the migrated code for\n"
        "single-unit review.\n"
        "• graphs/*\n"
        "call_graph + dataflow as DOT, Mermaid, and SVG\n"
        "(pre-rendered when graphviz is on PATH).\n"
        "• comprehension_cache.jsonl\n"
        "Content-hash keyed summaries for incremental reruns."
    ),
    "right_header": "Pilot run  ·  aws_legacy/generated_code  ·  4 Lambdas",
    "right_body": (
        "• Outcome: 4 / 4 succeeded, tests pass, 82 / 100\n"
        "Verify gate green on every unit; snapshot-rollback\n"
        "triggered once on a retry regression and recovered.\n"
        "• Cost: $2.86 of $5.00 across 684 LLM calls\n"
        "603k in / 136k out tokens  ·  gpt-5.4.\n"
        "Comprehension accounted for ~658 calls; translator\n"
        "+ reviewer + security the remainder.\n"
        "• Runtime: 17m 30s  (8m without comprehension)\n"
        "Extra time buys populated Purpose + 40+\n"
        "paragraph summaries per Lambda.\n"
        "• Cache artifact: 656 entries persisted\n"
        "Second run on unchanged source is near-free\n"
        "for the comprehension stage.\n"
        "• KG: 4 programs · 199 paragraphs\n"
        "40 external_call nodes  ·  5 dataset_ref nodes\n"
        "(DynamoDB, S3, SQS, SSM)."
    ),
    "footer": (
        "Part 3  ·  AWS Lambda pipeline  ·  4 Lambdas · $2.86 · 684 calls · 82/100 · "
        "artifacts persisted for review"
    ),
}


NEW_SLIDES = [
    # (insert_after_in_ORIGINAL_deck, content)
    #
    # The driver adjusts these for earlier inserts by adding
    # ``inserts_done`` to each, so quote original-deck indices here.
    # Sit both Part-2 additions right next to the Stage 3 LLM₁ slide
    # (original index 12) so they read as one logical block; the cache
    # slide is deliberately positioned AFTER large-file (so we insert
    # at original 12 twice — the second insert lands behind the first).
    (12, LARGE_FILE_SLIDE),
    (12, CACHE_KNOBS_SLIDE),
    # Last Part-3 slide in the original deck is 21 ("Reviewer rules");
    # our AWS pipeline slide closes out Part 3 before Part 4 opens.
    (21, AWS_PIPELINE_SLIDE),
]


# --------------------------------------------------------------------------- #
# Edits to existing slides (text rewrites only — no geometry changes).
# --------------------------------------------------------------------------- #

STAGE3_MECHANICS_UPDATE = (
    "• DFS post-order\n"
    "Leaves first; parents consume already-distilled child summaries.\n"
    "A 15K-LOC program never enters the prompt as raw source.\n"
    "• Auto-chunking for oversized nodes  (new)\n"
    "When a single node busts the input-token budget the summarizer\n"
    "splits at semantic boundaries, summarizes each chunk in parallel,\n"
    "then reduces with the node's kind prompt — 40k-LOC paragraphs\n"
    "go through cleanly instead of failing context.\n"
    "• Cached by node hash + chunk hash\n"
    "Re-running unchanged source short-circuits the LLM call.\n"
    "Cache lives on disk at <workdir>/comprehension_cache.jsonl\n"
    "and survives across runs + partial edits.\n"
    "• Cheap-model territory\n"
    "1–4k tokens / call, gpt-5-mini / haiku class.\n"
    "Budget tracker tags all calls as `comprehension`.\n"
    "• Output  →  KGNode.llm_summary\n"
    "What the node does, what it reads/writes, side effects.\n"
    "Plus an embedding for retrieval at synthesis time."
)


def _shape_by_name(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def apply_new_slide_content(slide, content: dict) -> None:
    """Map content dict to shapes on the cloned slide by position/role.

    The template (slide 12) has a fixed layout:
        TextBox 1  - header bar (PART X - ...)
        TextBox 2  - slide title
        TextBox 3  - subtitle
        TextBox 6  - pipeline step 1
        TextBox 9  - pipeline step 2
        TextBox 12 - pipeline step 3
        TextBox 15 - pipeline step 4
        TextBox 18 - pipeline step 5
        TextBox 21 - left-column header
        TextBox 22 - left-column body
        TextBox 25 - right-column header
        TextBox 26 - right-column body
        TextBox 27 - footer
    """
    mapping = {
        "TextBox 1":  content["title_bar"],
        "TextBox 2":  content["title"],
        "TextBox 3":  content["subtitle"],
        "TextBox 6":  content["pipe_1"],
        "TextBox 9":  content["pipe_2"],
        "TextBox 12": content["pipe_3"],
        "TextBox 15": content["pipe_4"],
        "TextBox 18": content["pipe_5"],
        "TextBox 21": content["left_header"],
        "TextBox 22": content["left_body"],
        "TextBox 25": content["right_header"],
        "TextBox 26": content["right_body"],
        "TextBox 27": content["footer"],
    }
    for shape_name, text in mapping.items():
        shape = _shape_by_name(slide, shape_name)
        if shape is None:
            print(f"  WARN: no shape named {shape_name} on this slide")
            continue
        _set_text(shape, text)


def update_stage3_slide(prs: Presentation) -> None:
    """Rewrite slide 12's left-column body to mention auto-chunking."""
    slide = prs.slides[12]
    shape = _shape_by_name(slide, "TextBox 22")
    if shape is None:
        return
    _set_text(shape, STAGE3_MECHANICS_UPDATE)


def update_pilot_results_slide(prs: Presentation) -> None:
    """Tweak Part 4 pilot slide footer to acknowledge the Part 3 pilot.

    After the three inserts (at positions 13, 14, 24), the original COBOL
    pilot slide (was index 29) now sits at index 32. Editing it updates
    the footer only; the stats + headers stay untouched.
    """
    slide = prs.slides[32]
    # The pilot slide's footer is TextBox 22 per the original layout
    # dump — a single-line bottom caption.
    footer = _shape_by_name(slide, "TextBox 22")
    if footer is None:
        return
    _set_text(
        footer,
        "Part 4  ·  COBOL pilot  ·  see AWS Lambda pilot in Part 3  ·  "
        "full breakdown in pilot/PILOT_REPORT.md",
    )


# --------------------------------------------------------------------------- #
# Driver
# --------------------------------------------------------------------------- #


def main() -> None:
    prs = Presentation(str(DECK))
    orig_count = len(prs.slides)
    print(f"loaded deck with {orig_count} slides")

    # Each clone_slide() appends to the end; reorder_slide() moves it into
    # place. Because earlier inserts shift every later index by +1, we
    # supply insert-after positions in the ORIGINAL ordering and apply them
    # in top-down order, adjusting by the number of inserts already made.
    inserts_done = 0
    for original_insert_after, content in NEW_SLIDES:
        insert_after = original_insert_after + inserts_done
        # Template is always slide 12 — unchanged by earlier inserts
        # because they all landed AFTER position 12 (or at position 12+,
        # never before). Cloning uses the constant index rather than a
        # shifted one.
        new_slide = clone_slide(prs, TEMPLATE_SLIDE_INDEX)
        apply_new_slide_content(new_slide, content)
        reorder_slide(prs, new_slide, insert_after)
        inserts_done += 1
        print(f"inserted new slide after index {insert_after} ({content['title']})")

    update_stage3_slide(prs)
    print("updated slide 12 (Stage 3 LLM1 mechanics)")

    update_pilot_results_slide(prs)
    print("updated slide 32 (COBOL pilot results footer)")

    prs.save(str(DECK))
    print(f"saved {DECK}  (slides: {orig_count} → {len(prs.slides)})")


if __name__ == "__main__":
    main()
