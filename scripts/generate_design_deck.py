"""Build the legacy-modernization platform design deck.

Outputs ``docs/design_deck.pptx``. The deck is structured into five parts:

  Part 0 — Cover
  Part 1 — The legacy-modernization problem
  Part 2 — Platform architecture (generic, cartridge-driven)
  Part 3 — Sample cartridge: AWS Lambda → Azure Functions
  Part 4 — Sample cartridge: COBOL + JCL + DB2 + CICS → Java 25 / Spring Boot 4.0.5
  Part 5 — Validation, results, roadmap

Produced with ``python-pptx``. The output is 16:9, uses only standard
shapes/colors so it imports cleanly into Google Slides.

Run from repo root:

    python scripts/generate_design_deck.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = REPO_ROOT / "docs" / "design_deck.pptx"

# ---------------------------------------------------------------------- #
# Colour palette — kept conservative so it reads well in Google Slides.   #
# ---------------------------------------------------------------------- #
PRIMARY = RGBColor(0x0F, 0x3E, 0x6E)   # deep blue — titles, headers
ACCENT = RGBColor(0x1E, 0x7A, 0xB8)    # lighter blue — dividers, callouts
DARK = RGBColor(0x1A, 0x1A, 0x1A)      # body text
GRAY = RGBColor(0x6B, 0x72, 0x80)      # secondary text
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)     # card background
BORDER = RGBColor(0xD1, 0xD5, 0xDB)
CODE_BG = RGBColor(0x1F, 0x29, 0x37)
CODE_FG = RGBColor(0xE5, 0xE7, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x74, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
TITLE_TOP = Inches(0.35)
TITLE_HEIGHT = Inches(0.7)
SUBTITLE_TOP = Inches(1.05)
SUBTITLE_HEIGHT = Inches(0.35)
BODY_TOP = Inches(1.55)
FOOTER_TOP = Inches(7.05)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------- #
# Helpers                                                                 #
# ---------------------------------------------------------------------- #


def _set_text(tf, runs, *, size=Pt(14), color=DARK, bold=False, mono=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, clear=True):
    """Write one or more runs into a text frame.

    ``runs`` is a list of str OR a list of ``(text, style_dict)``. Each run
    gets its own paragraph (one-paragraph-per-run keeps Google Slides happy).
    """
    if clear:
        tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    items = runs if isinstance(runs, (list, tuple)) else [runs]
    for item in items:
        if isinstance(item, tuple):
            text, style = item
        else:
            text, style = item, {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = style.get("align", align)
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = style.get("size", size)
        f.bold = style.get("bold", bold)
        f.italic = style.get("italic", False)
        f.name = "Menlo" if style.get("mono", mono) else "Calibri"
        f.color.rgb = style.get("color", color)


def add_rect(slide, left, top, width, height, *,
             fill=None, line=None, line_weight=Pt(0.75),
             text=None, text_kwargs=None,
             shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_weight
    if text is not None:
        _set_text(shp.text_frame, text, **(text_kwargs or {}))
    return shp


def add_textbox(slide, left, top, width, height, text, **kwargs):
    tb = slide.shapes.add_textbox(left, top, width, height)
    _set_text(tb.text_frame, text, **kwargs)
    return tb


def add_title(slide, title_text, subtitle_text=None, *, section=None):
    # Section label (e.g. "PART 2 — PLATFORM ARCHITECTURE")
    if section is not None:
        add_textbox(
            slide, MARGIN, Inches(0.25), Inches(12), Inches(0.3),
            [(section, {"size": Pt(11), "color": ACCENT, "bold": True})],
        )
    add_textbox(
        slide, MARGIN, TITLE_TOP, Inches(12), TITLE_HEIGHT,
        [(title_text, {"size": Pt(30), "color": PRIMARY, "bold": True})],
    )
    if subtitle_text:
        add_textbox(
            slide, MARGIN, SUBTITLE_TOP, Inches(12), SUBTITLE_HEIGHT,
            [(subtitle_text, {"size": Pt(14), "color": GRAY})],
        )
    # Accent underline
    add_rect(
        slide, MARGIN, Inches(1.45), Inches(0.9), Emu(38100),  # ~3 pt
        fill=ACCENT,
    )


def add_footer(slide, text):
    add_textbox(
        slide, MARGIN, FOOTER_TOP, Inches(12), Inches(0.3),
        [(text, {"size": Pt(9), "color": GRAY})],
    )


def blank_slide():
    return prs.slides.add_slide(BLANK)


def bullet_paragraphs(tf, bullets, *, size=Pt(14), color=DARK,
                      sub_size=Pt(12), sub_color=GRAY):
    """Render a nested-bullet list.

    ``bullets`` is a list of items. Each item is:
      - str → single top-level bullet
      - (header, [child_strings]) → header bullet plus nested lines
      - (header, [child_strings], style_override_dict)
    """
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, str):
            header, children = item, []
        else:
            header, children = item[0], item[1]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        r = p.add_run()
        r.text = "• " + header
        r.font.size = size
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = "Calibri"
        for child in children:
            if isinstance(child, tuple):
                ctext, cstyle = child
            else:
                ctext, cstyle = child, {}
            cp = tf.add_paragraph()
            cp.level = 1
            cr = cp.add_run()
            cr.text = ctext
            cr.font.size = cstyle.get("size", sub_size)
            cr.font.bold = cstyle.get("bold", False)
            cr.font.italic = cstyle.get("italic", False)
            cr.font.color.rgb = cstyle.get("color", sub_color)
            cr.font.name = "Menlo" if cstyle.get("mono", False) else "Calibri"


def code_block(slide, left, top, width, height, code, *, font_size=Pt(11)):
    add_rect(
        slide, left, top, width, height,
        fill=CODE_BG, line=CODE_BG,
    )
    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = code.splitlines() or [""]
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.size = font_size
        r.font.name = "Menlo"
        r.font.color.rgb = CODE_FG


def card(slide, left, top, width, height, title, lines, *, accent=ACCENT):
    """Rounded card: colored title band + body bullets."""
    add_rect(
        slide, left, top, width, height,
        fill=LIGHT, line=BORDER,
    )
    add_rect(
        slide, left, top, width, Inches(0.45),
        fill=accent, line=accent,
        text=[(title, {"size": Pt(13), "color": RGBColor(0xFF, 0xFF, 0xFF),
                       "bold": True, "align": PP_ALIGN.LEFT})],
        text_kwargs={"anchor": MSO_ANCHOR.MIDDLE},
    )
    body_tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.55),
        width - Inches(0.3), height - Inches(0.65),
    )
    bullet_paragraphs(body_tb.text_frame, lines, size=Pt(11), sub_size=Pt(10))


def two_column(slide, *, left_title, left_items, right_title, right_items,
               top=Inches(1.7), height=Inches(5.1)):
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    card(slide, MARGIN, top, col_w, height, left_title, left_items)
    card(slide, MARGIN + col_w + Inches(0.3), top, col_w, height,
         right_title, right_items, accent=PRIMARY)


def pipeline_boxes(slide, top, height, boxes, *, arrow_gap=Inches(0.15)):
    """Horizontal row of boxes with connecting arrows."""
    n = len(boxes)
    total_w = SLIDE_W - 2 * MARGIN
    arrow_w = Inches(0.25)
    total_arrows = (n - 1) * (arrow_w + 2 * arrow_gap)
    box_w = (total_w - total_arrows) / n
    x = MARGIN
    for i, (title, sub) in enumerate(boxes):
        add_rect(
            slide, x, top, box_w, height,
            fill=LIGHT, line=ACCENT, line_weight=Pt(1.25),
            text=[
                (title, {"size": Pt(14), "color": PRIMARY, "bold": True,
                         "align": PP_ALIGN.CENTER}),
                (sub, {"size": Pt(10), "color": GRAY,
                       "align": PP_ALIGN.CENTER}),
            ],
            text_kwargs={"anchor": MSO_ANCHOR.MIDDLE},
        )
        x += box_w
        if i < n - 1:
            x += arrow_gap
            add_rect(
                slide, x, top + height / 2 - Inches(0.08), arrow_w, Inches(0.16),
                fill=ACCENT, line=ACCENT,
                shape=MSO_SHAPE.RIGHT_ARROW,
            )
            x += arrow_w + arrow_gap


def add_table(slide, left, top, width, height, headers, rows,
              *, header_fill=PRIMARY, header_fg=RGBColor(0xFF, 0xFF, 0xFF),
              body_font=Pt(11), header_font=Pt(12)):
    cols = len(headers)
    rows_count = len(rows) + 1
    tbl = slide.shapes.add_table(rows_count, cols, left, top, width, height).table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h
        r.font.size = header_font
        r.font.bold = True
        r.font.color.rgb = header_fg
        r.font.name = "Calibri"
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = body_font
            r.font.name = "Menlo" if ci == 0 else "Calibri"
            r.font.color.rgb = DARK
    return tbl


# ---------------------------------------------------------------------- #
# Slide builders                                                          #
# ---------------------------------------------------------------------- #


def s00_cover():
    s = blank_slide()
    # Full-bleed gradient-ish band
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(3.8), fill=PRIMARY, line=PRIMARY)
    add_rect(s, Inches(0), Inches(3.8), SLIDE_W, Inches(0.12), fill=ACCENT, line=ACCENT)
    add_textbox(
        s, MARGIN, Inches(1.0), Inches(12), Inches(0.5),
        [("LEGACY MODERNIZATION PLATFORM",
          {"size": Pt(14), "color": RGBColor(0xBB, 0xD0, 0xE6), "bold": True})],
    )
    add_textbox(
        s, MARGIN, Inches(1.6), Inches(12), Inches(1.4),
        [("A cartridge-driven framework for",
          {"size": Pt(32), "color": RGBColor(0xFF, 0xFF, 0xFF), "bold": True}),
         ("LLM-assisted migration of legacy estates",
          {"size": Pt(32), "color": RGBColor(0xFF, 0xFF, 0xFF), "bold": True})],
    )
    add_textbox(
        s, MARGIN, Inches(4.3), Inches(12), Inches(0.4),
        [("Generic platform design  ·  Two sample cartridges",
          {"size": Pt(18), "color": PRIMARY, "bold": True})],
    )
    add_textbox(
        s, MARGIN, Inches(4.85), Inches(12), Inches(2),
        [("1.  AWS Lambda (Python / Node / TS / Java / C#) → Azure Functions (Python)",
          {"size": Pt(14), "color": DARK}),
         ("2.  COBOL + JCL + DB2 + CICS → Java 25 LTS / Spring Boot 4.0.5",
          {"size": Pt(14), "color": DARK}),
         ("", {"size": Pt(6)}),
         ("Built on Microsoft Agent Framework  ·  Multi-provider LLM  ·  Budget-gated",
          {"size": Pt(12), "color": GRAY, "italic": True})],
    )
    add_footer(s, "Legacy Modernization Platform  ·  Design Deck")


def s01_contents():
    s = blank_slide()
    add_title(s, "What's in this deck",
              subtitle_text="Five parts. Read sequentially, or jump to the sample that matches your estate.")
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    card(s, MARGIN, Inches(1.8), col_w, Inches(5),
         "Generic platform",
         [
             ("Part 1 — The legacy-modernization problem + influences", [
                 "Fowler, Thoughtworks, Anthropic — what we took",
                 "from each (full references at the end)",
             ]),
             ("Part 2 — Platform architecture", [
                 "Cartridge model, Knowledge Graph, comprehension, planning,",
                 "worker loop, agents, budget, cross-unit API injection",
             ]),
             "Part 5 — Validation, results, roadmap, references",
         ])
    card(s, MARGIN + col_w + Inches(0.3), Inches(1.8), col_w, Inches(5),
         "Two sample cartridges",
         [
             ("Part 3 — AWS Lambda → Azure Functions", [
                 "Polyglot source (5 languages) → Python Azure Functions",
                 "Deterministic recipes + residuals-driven LLM",
             ]),
             ("Part 4 — COBOL + JCL + DB2 + CICS → Java 25 / Spring Boot 4.0.5", [
                 "Target-style classifier (batch / CICS / subroutine)",
                 "Cross-unit API injection, data-dictionary rendering,",
                 "dual-run equivalence harness",
             ]),
         ], accent=PRIMARY)
    add_footer(s, "Contents")


def s02_problem():
    s = blank_slide()
    add_title(s, "Why modernizing legacy estates is hard",
              subtitle_text="The artefacts you're migrating are alive, entangled, and under-documented.")
    two_column(
        s,
        left_title="The estate problem",
        left_items=[
            ("Scale & heterogeneity", [
                "Tens of thousands of artefacts per enterprise",
                "Multiple languages, stacks, dialects, tribal extensions",
            ]),
            ("Coupling you can't see", [
                "Implicit contracts (COMMAREA, COMP-3, DD cards)",
                "Shared copybooks, JCL step ordering, CICS transaction maps",
            ]),
            ("Lost knowledge", [
                "Authors have retired; specs don't exist",
                "Only the code tells the truth",
            ]),
            ("Correctness bar", [
                "Dual-run against prod inputs is the acceptance criterion,",
                "not \"looks right\" to a reviewer",
            ]),
        ],
        right_title="Why LLMs alone aren't enough",
        right_items=[
            ("Context blindness", [
                "LLM sees one file at a time; misses cross-program coupling",
                "Hallucinates DTO names, method signatures, dependencies",
            ]),
            ("No deterministic baseline", [
                "Same input → different output run-to-run",
                "Hard to diff, re-review, or attribute regressions",
            ]),
            ("No correctness gate", [
                "Code that ``compiles`` is not code that ``does the same thing``",
            ]),
            ("Cost unpredictability", [
                "One naive prompt × 50K files = five-figure bill",
                "No way to cap or halt mid-run",
            ]),
        ],
    )
    add_footer(s, "Part 1  ·  The legacy-modernization problem")


def s03_thesis():
    s = blank_slide()
    add_title(s, "Our thesis",
              subtitle_text="Build a platform, not a script. Treat migration like a compiler pipeline — with LLMs at the stages that genuinely need reasoning.  Influenced directly by Fowler / Thoughtworks / Anthropic (see references).")
    # Three-card layout
    w = (SLIDE_W - 2 * MARGIN - 2 * Inches(0.3)) / 3
    h = Inches(5)
    top = Inches(1.8)
    card(s, MARGIN, top, w, h,
         "Deterministic first",
         [
             ("Static analysis over guessing", [
                 "Parsers + knowledge graph + CRUD matrix do the",
                 "heavy lifting before any LLM call is made",
             ]),
             ("Recipes before prompts", [
                 "libcst / tree-sitter transforms handle mechanical rewrites",
                 "LLM only sees what recipes couldn't fix (``residuals``)",
             ]),
         ], accent=ACCENT)
    card(s, MARGIN + w + Inches(0.3), top, w, h,
         "LLM where it matters",
         [
             ("Translation, review, security", [
                 "Agents specialise per role; prompts are small and precise",
             ]),
             ("Retry + snapshot-rollback", [
                 "Every attempt is committed to a snapshot;",
                 "regressions roll back to the last known-good state",
             ]),
             ("Cross-unit API injection", [
                 "Callee's already-migrated Java is inlined into the",
                 "caller's prompt verbatim — no hallucinated signatures",
             ]),
         ], accent=PRIMARY)
    card(s, MARGIN + 2 * w + 2 * Inches(0.3), top, w, h,
         "Operationally safe",
         [
             ("Budget-gated", [
                 "``max_budget_usd`` per run; per-call accounting;",
                 "halts cleanly when the cap is reached",
             ]),
             ("Checkpointed", [
                 "Every wave commits a checkpoint — resume or post-mortem",
                 "without re-running successful units",
             ]),
             ("Cartridge-isolated", [
                 "New source / target pair = new cartridge directory.",
                 "Platform core doesn't change.",
             ]),
         ], accent=GREEN)
    add_footer(s, "Part 1  ·  Platform thesis")


def s04_influences():
    s = blank_slide()
    add_title(s, "Influences: what we took from Fowler / Thoughtworks / Anthropic",
              subtitle_text="The platform isn't invented from scratch.  These four references shape the core pipeline, the comprehension stage, and the cost narrative.")
    # Full-width table mapping reference → pattern adopted
    add_table(
        s, MARGIN, Inches(1.85), SLIDE_W - 2 * MARGIN, Inches(4.6),
        ["Reference", "Pattern we adopted", "Where it lives in our code"],
        [
            ("Martin Fowler — Legacy Modernization meets Generative AI\n(martinfowler.com, Ryan et al.)",
             "CodeConcise-style pipeline: AST → Knowledge Graph → post-order\nDFS summaries → business spec. CRUD matrix + Feathers-style seams\nover the KG. Dual-run equivalence as the acceptance bar.",
             "platform_core/kg/  ·  comprehension/summarizer.py\ncomprehension/business_spec.py  ·  pipeline/crud.py\nplatform_core/dualrun/  ·  kg/render.py (DOT + Mermaid)"),
            ("Thoughtworks — Claude Code for COBOL Modernization: a reality check\n(thoughtworks.com, 2025)",
             "Practical findings from running agentic workflows on real mainframe\ncorpora: structured specs beat raw code, human-in-loop at review,\ntarget-style differentiation per program, compile gate is non-negotiable.",
             "cartridges/cobol_to_java25_springboot/\n  target_style.py (batch / cics / subroutine split)\n  prompts/translator_*.md  ·  prompts/reviewer.md\n  verify_unit() in cartridge.py (mvn compile gate)"),
            ("Anthropic — How AI helps break the cost barrier for COBOL modernization\n(claude.com/blog)",
             "Agentic retry loops over structured JSON verdicts; bounded-budget\nmulti-agent orchestration; caller/callee ordering so the LLM sees\nthe real API instead of hallucinating one.",
             "platform_core/maf_workflows/budget.py  (per-call cost + cap)\nunit_worker.py  (retry loop + snapshot rollback)\ncartridge._collect_migrated_callee_apis()\n  (cross-unit API injection)"),
            ("Broader Fowler canon — Working Effectively with Legacy Code (Feathers),\nRefactoring (Fowler)",
             "Seams as refactoring pivots; characterisation-test mindset →\ndual-run fixtures; avoid-premature-abstraction → cartridge isolation.",
             "platform_core/pipeline/crud.py  (seam ranking)\nplatform_core/dualrun/fixtures.py  ·  verdict.py\nplatform_core/cartridge/  (narrow hook surface)"),
        ],
        body_font=Pt(10), header_font=Pt(11),
    )
    add_textbox(
        s, MARGIN, Inches(6.55), SLIDE_W - 2 * MARGIN, Inches(0.4),
        [("Explicit citations are on the final References slide; every URL is the source of record.",
          {"size": Pt(11), "color": GRAY, "italic": True,
           "align": PP_ALIGN.CENTER})],
    )
    add_footer(s, "Part 1  ·  Influences")


def s10_architecture():
    s = blank_slide()
    add_title(s, "Platform architecture at a glance",
              subtitle_text="One generic pipeline. Cartridges plug into every stage without touching core code.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    # Pipeline bar across the middle
    pipeline_top = Inches(2.3)
    pipeline_boxes(
        s, pipeline_top, Inches(1.0),
        [
            ("Discover", "walk repo, parse adapters,\nclassify units"),
            ("Graph", "build cross-unit KG,\nemit grapher_hints"),
            ("Plan", "topological waves\nfrom KG edges"),
            ("Worker", "recipes → translator\n→ reviewer → verify"),
            ("Verify", "mvn compile +\ndual-run gate"),
        ],
    )
    # Underneath: cartridge interface as the glue
    add_rect(s, MARGIN, Inches(4.0), SLIDE_W - 2 * MARGIN, Inches(0.45),
             fill=ACCENT, line=ACCENT,
             text=[("Cartridge interface (plugin boundary)",
                    {"size": Pt(13), "bold": True,
                     "color": RGBColor(0xFF, 0xFF, 0xFF),
                     "align": PP_ALIGN.CENTER})],
             text_kwargs={"anchor": MSO_ANCHOR.MIDDLE})
    # Cartridge hooks laid out horizontally
    hooks = [
        ("adapters()", "parse source\n→ KG nodes"),
        ("unit_classifier()", "files → units"),
        ("grapher_hints()", "cross-unit edges"),
        ("recipes()", "deterministic\ntransforms"),
        ("build_translator_request()", "per-unit prompt\n+ callee injection"),
        ("verify_unit()", "compile + dual-run"),
    ]
    w_h = (SLIDE_W - 2 * MARGIN) / len(hooks)
    for i, (name, desc) in enumerate(hooks):
        x = MARGIN + w_h * i
        add_rect(s, x + Inches(0.05), Inches(4.55), w_h - Inches(0.1), Inches(1.0),
                 fill=LIGHT, line=BORDER,
                 text=[(name, {"size": Pt(11), "bold": True, "color": PRIMARY,
                               "align": PP_ALIGN.CENTER, "mono": True}),
                       (desc, {"size": Pt(9), "color": GRAY,
                               "align": PP_ALIGN.CENTER})],
                 text_kwargs={"anchor": MSO_ANCHOR.MIDDLE})
    # Stacks: observability + ops
    add_rect(s, MARGIN, Inches(5.85), SLIDE_W - 2 * MARGIN, Inches(1.0),
             fill=PRIMARY, line=PRIMARY,
             text=[("Cross-cutting:  budget tracker  ·  checkpoints  ·  snapshot-rollback  ·  per-role agent factory  ·  retry loop",
                    {"size": Pt(12), "color": RGBColor(0xFF, 0xFF, 0xFF),
                     "align": PP_ALIGN.CENTER, "bold": True})],
             text_kwargs={"anchor": MSO_ANCHOR.MIDDLE})
    add_footer(s, "Part 2  ·  Architecture overview")


def s11_cartridge_model():
    s = blank_slide()
    add_title(s, "The cartridge: one migration family per plugin",
              subtitle_text="A ``MigrationCartridge`` is a Python class that exposes the hooks the pipeline calls. Nothing else in the codebase needs to change.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    code = """from maf_generic_migrator_v1.platform_core.cartridge import (
    AgentSpec, MigrationCartridge,
)

class CobolToJava25SpringBoot(MigrationCartridge):
    id = "cobol_to_java25_springboot"
    description = "COBOL (batch + CICS + DB2) -> Java 25 LTS / Spring Boot 4.0.5"

    def adapters(self):                  # source parsers
        return {"cobol": CobolAdapter()}

    def unit_classifier(self, repo_root): ...   # files  -> units
    def grapher_hints(self, inventory):   ...   # CALL edges  -> KG
    def recipes(self):                    ...   # deterministic transforms
    def translator_agents(self): ...             # 3 specs (batch/cics/subroutine)
    def reviewer_agents(self):  ...
    async def build_translator_request(   ...   # per-unit prompt
        self, item, workdir, cfg, agent_bundles, ...
    ):
    def verify_unit(self, workdir, cfg):  ...   # mvn compile + dual-run"""
    code_block(s, MARGIN, Inches(1.85), Inches(8.2), Inches(5.0), code)
    card(s, Inches(9.0), Inches(1.85), Inches(3.8), Inches(5.0),
         "Why this shape",
         [
             ("Narrow surface area", [
                 "<20 hooks; most are optional",
             ]),
             ("Conformance-tested", [
                 "Platform ships a test suite every cartridge must pass",
             ]),
             ("Drop-in and go", [
                 "New cartridge = one directory;",
                 "``.env CARTRIDGE=<id>`` picks it",
             ]),
             ("Isolated blast radius", [
                 "A bug in one cartridge can't break others",
             ]),
         ], accent=GREEN)
    add_footer(s, "Part 2  ·  Cartridge model")


def s12_discovery():
    s = blank_slide()
    add_title(s, "Discovery: turn a repo into an inventory",
              subtitle_text="Language adapters walk the tree, parse per-file, and emit KG nodes + unit IDs.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    two_column(
        s,
        left_title="Language adapter contract",
        left_items=[
            ("One adapter per source language", [
                "`.cbl` / `.cpy` → CobolAdapter",
                "`.py` / `.ts` / `.java` / `.cs` → language-specific adapters",
            ]),
            ("Parsing strategy", [
                "tree-sitter for structured AST where a grammar exists",
                "Regex fallback for legacy formats (EXEC SQL / CICS / JCL)",
            ]),
            ("Output: KG nodes + edges", [
                "Program, Section, Paragraph, Field, Record, File,",
                "SqlBlock, TxnCall, ExternalCall, Dataset, Job, Step",
            ]),
            ("Unit classification", [
                "A ``unit`` is the coherent translation target",
                "(typically one Maven module per COBOL program,",
                "one Azure Function per Lambda handler)",
            ]),
        ],
        right_title="What discovery produces",
        right_items=[
            ("Inventory IR", [
                "list of ``UnitIR`` objects with files, imports,",
                "language, handler entry point",
            ]),
            ("Populated Knowledge Graph", [
                "All nodes + intra-unit edges (contains, references)",
            ]),
            ("Warnings surface", [
                "Files the adapter couldn't parse flagged early —",
                "never a silent drop",
            ]),
            ("Ready for the grapher", [
                "Cartridge-specific ``grapher_hints()`` adds the",
                "cross-unit edges (CALL graph, Lambda invocations,",
                "SQL FK joins)",
            ]),
        ],
    )
    add_footer(s, "Part 2  ·  Discovery")


def s13_knowledge_graph():
    s = blank_slide()
    add_title(s, "The Knowledge Graph: one queryable model of the estate",
              subtitle_text="Typed nodes + typed edges. Backed by NetworkX; rendered as DOT or Mermaid on demand.  Schema follows Fowler's CodeConcise pattern.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    # Left: node kinds table. Right: edge kinds table.
    tbl_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    tbl_top = Inches(1.85)
    tbl_h = Inches(4.9)
    add_textbox(
        s, MARGIN, tbl_top - Inches(0.35), tbl_w, Inches(0.3),
        [("Node kinds  (structural artefacts)",
          {"size": Pt(13), "bold": True, "color": PRIMARY})],
    )
    add_table(
        s, MARGIN, tbl_top, tbl_w, tbl_h,
        ["Kind", "Represents"],
        [
            ("program",       "One COBOL program / Lambda handler / etc."),
            ("section",       "COBOL SECTION or logical block"),
            ("paragraph",     "COBOL paragraph / function / method"),
            ("field",         "Single data field (with PIC / type)"),
            ("record",        "FD / 01-level record; group item"),
            ("file",          "Opened file (FD) or storage ref"),
            ("sql_block",     "One EXEC SQL statement"),
            ("txn_call",      "EXEC CICS / transaction invocation"),
            ("external_call", "CALL 'TARGET' / cross-module invocation"),
            ("dataset_ref",   "JCL DD → physical dataset"),
            ("job / step",    "JCL job and step"),
        ],
        body_font=Pt(10), header_font=Pt(11),
    )
    add_textbox(
        s, MARGIN + tbl_w + Inches(0.3), tbl_top - Inches(0.35), tbl_w, Inches(0.3),
        [("Edge kinds  (relationships)",
          {"size": Pt(13), "bold": True, "color": PRIMARY})],
    )
    add_table(
        s, MARGIN + tbl_w + Inches(0.3), tbl_top, tbl_w, tbl_h,
        ["Kind", "Connects"],
        [
            ("contains",     "parent → child (program → paragraph → field)"),
            ("references",   "field → field (MOVE, COMPUTE)"),
            ("reads",        "program → sql_block / file  (CRUD: R)"),
            ("writes",       "program → sql_block / file  (CRUD: W)"),
            ("calls",        "caller program → callee program"),
            ("touches",      "job step → dataset  (JCL binding)"),
            ("step_of",      "step → job  (JCL ordering)"),
            ("redefines",    "field → field  (COBOL REDEFINES)"),
            ("input_of",     "field → program  (LINKAGE, inbound)"),
            ("output_of",    "field → program  (LINKAGE, outbound)"),
        ],
        body_font=Pt(10), header_font=Pt(11),
    )
    add_footer(s, "Part 2  ·  Knowledge Graph schema")


def s14_comprehension():
    s = blank_slide()
    add_title(s, "Comprehension: give the LLM structured understanding, not raw code",
              subtitle_text="Post-order DFS summaries + CRUD matrix + business spec = prompt payload that beats ``paste the whole file``.  Directly adapted from the Fowler / Thoughtworks CodeConcise pipeline.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    # Flow pipeline
    pipeline_boxes(
        s, Inches(1.85), Inches(0.9),
        [
            ("KG",           "nodes + edges\nfrom discovery"),
            ("Summariser",   "DFS post-order\nper-paragraph"),
            ("CRUD Matrix",  "program × data\n→ R / W / RW"),
            ("Seam Ranking", "coupling score\n(Thoughtworks)"),
            ("Business Spec","Purpose / IO /\nSide effects"),
        ],
    )
    # Two cards underneath
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    top = Inches(3.15)
    h = Inches(3.75)
    card(s, MARGIN, top, col_w, h,
         "Summariser",
         [
             ("DFS post-order", [
                 "Leaves first; parents consume child summaries",
             ]),
             ("Cached by node hash", [
                 "Re-running doesn't re-summarise unchanged paragraphs",
             ]),
             ("Output", [
                 "One paragraph per KG node: what it does, what it",
                 "reads/writes, its side effects",
             ]),
         ])
    card(s, MARGIN + col_w + Inches(0.3), top, col_w, h,
         "Business spec (per unit)",
         [
             ("Purpose / Inputs / Outputs / Side effects / Invariants", [
                 "Derived from KG + summaries; no hand-writing",
             ]),
             ("Target-style aware", [
                 "Batch units list datasets; subroutines list LINKAGE fields",
                 "CICS units list transaction maps + COMMAREA",
             ]),
             ("Feeds the translator prompt", [
                 "Not the raw COBOL — the interpreted intent",
             ]),
         ], accent=PRIMARY)
    add_footer(s, "Part 2  ·  Comprehension pipeline")


def s15_planning():
    s = blank_slide()
    add_title(s, "Planning: topological waves from the Knowledge Graph",
              subtitle_text="Dependencies in the KG drive parallelism. Callees ship before callers.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    two_column(
        s,
        left_title="How waves are computed",
        left_items=[
            ("Input: cross-unit edges", [
                "``grapher_hints()`` emits ``calls`` / ``references``",
                "edges between UnitIR objects",
            ]),
            ("Kahn's algorithm", [
                "Top-sort the unit DAG into waves",
                "Each wave = units whose dependencies all shipped",
            ]),
            ("Parallelism knob", [
                "``MAX_PARALLEL`` bounds fan-out per wave",
                "Cheap units batch, expensive ones serialise",
            ]),
            ("Ambiguity handling", [
                "``calls`` edges to units outside the inventory are dropped",
                "Reported but don't block the plan",
            ]),
        ],
        right_title="Why this matters for LLM quality",
        right_items=[
            ("Callees shipped first", [
                "When STMTGEN's translator runs, INTCALC's Java is",
                "already on disk — we inline its real API into",
                "STMTGEN's prompt (see cross-unit API injection)",
            ]),
            ("No hallucinated signatures", [
                "Previous runs emitted stub classes with invented method",
                "names; now the real signature is in the prompt verbatim",
            ]),
            ("Early failures halt cheap", [
                "A callee that fails compilation stops its callers",
                "from being attempted",
            ]),
            ("Deterministic ordering", [
                "Same KG  →  same wave layout;",
                "reproducibility for regression triage",
            ]),
        ],
    )
    add_footer(s, "Part 2  ·  Planner")


def s16_worker():
    s = blank_slide()
    add_title(s, "Unit worker: the per-unit retry loop",
              subtitle_text="Same loop for every cartridge. Bounded retries, snapshot-rollback on regression.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    # Flow diagram + annotations
    top = Inches(1.85)
    h = Inches(0.85)
    pipeline_boxes(
        s, top, h,
        [
            ("Recipes",     "deterministic\ntransforms"),
            ("Residuals",   "what recipes\ncouldn't fix"),
            ("Translator",  "LLM agent\ngenerates files"),
            ("Parser",      "fenced blocks\n→ files on disk"),
            ("Verify",      "mvn compile +\ndual-run"),
            ("Reviewer",    "JSON verdict:\naccept/revise/reject"),
        ],
    )
    card(s, MARGIN, Inches(3.3), SLIDE_W - 2 * MARGIN, Inches(3.6),
         "Retry loop invariants",
         [
             ("Snapshot on every verify=pass", [
                 "Workdir is frozen; next attempt can be rolled back to it",
             ]),
             ("Reviewer downgrades to `revise` when verify passes", [
                 "``The compiler is authoritative.`` A reviewer phantom-rejecting",
                 "a file it literally can see in the workdir listing is ignored.",
             ]),
             ("Retry budget", [
                 "``max_reviewer_retries`` (default 2); merge reviewer + verify",
                 "feedback into the next prompt",
             ]),
             ("Regression triggers rollback", [
                 "If attempt N compiled but attempt N+1 doesn't, restore the",
                 "snapshot from attempt N and accept the last-known-good output",
             ]),
             ("Fail-hard guard rails", [
                 "No ``pom.xml`` emitted after parsing → fail the unit,",
                 "don't pretend success.  The parser accepts 5 path-header",
                 "conventions (``path=`` inline, ``<!-- -->``, ``//``, ``#``, ``/*  */``).",
             ]),
         ], accent=ACCENT)
    add_footer(s, "Part 2  ·  Unit worker loop")


def s17_agents():
    s = blank_slide()
    add_title(s, "Agents: role-specific LLMs with bounded responsibilities",
              subtitle_text="One chat client per role. Role-keyed ``max_tokens`` so nothing gets truncated mid-class.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    add_table(
        s, MARGIN, Inches(1.9), SLIDE_W - 2 * MARGIN, Inches(3.0),
        ["Role", "What it does", "max_tokens", "Output shape"],
        [
            ("translator-*", "Emit the full target module (one prompt per target style: batch / cics / subroutine / lambda).",
             "16000", "fenced code blocks with path= headers"),
            ("reviewer",    "JSON verdict (accept / revise / reject) + rule-slugged findings. Cartridge supplies a rule catalog in reviewer.md.",
             "8000",  "single JSON object"),
            ("security",    "OWASP + secrets scan over emitted code. Read-only.",
             "4000",  "findings list"),
            ("tester",      "Execute tests (when runner is wired) and parse results.",
             "4000",  "results summary"),
        ],
        body_font=Pt(10),
    )
    card(s, MARGIN, Inches(5.05), SLIDE_W - 2 * MARGIN, Inches(1.75),
         "Why separate agents",
         [
             ("Different prompts, different temperature, different budgets", [
                 "A reviewer on 16k tokens over-explains; a translator on 1k truncates.",
             ]),
             ("Failure isolation", [
                 "One agent crashing doesn't poison the others; result carries ``error`` and run continues.",
             ]),
             ("Composability", [
                 "Swap OpenAI → Anthropic → Ollama via ``LEGACY_MOD_PROVIDER``; agent code unchanged.",
             ]),
         ], accent=PRIMARY)
    add_footer(s, "Part 2  ·  Agent factory")


def s18_budget_checkpoint():
    s = blank_slide()
    add_title(s, "Budget & checkpoints: you can stop the run and know what it cost",
              subtitle_text="Per-call token accounting + per-run USD cap + resumable checkpoints.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    two_column(
        s,
        left_title="Budget tracker",
        left_items=[
            ("Per-call accounting", [
                "Every ``agent.run()`` reports input/output tokens",
                "Tracker maps model → (USD_in / 1M tok, USD_out / 1M tok)",
            ]),
            ("Halt on cap", [
                "``MAX_BUDGET_USD`` = hard cap; subsequent calls short-circuit",
                "Unit result carries ``budget exhausted before <tag>``",
            ]),
            ("By-model / by-role breakdown", [
                "Final log:  ``gpt-5.4: $0.4980 (17 calls, 90k in / 27k out)``",
                "Per-role counts:  translator=3, reviewer=7, security=3",
            ]),
            ("Cost report on disk", [
                "``cost_reports/<run_id>.json`` + append-only",
                "``cost_log.jsonl``  (one line per run)",
            ]),
        ],
        right_title="Checkpoints",
        right_items=[
            ("One per wave completion", [
                "``checkpoints/<run_id>.json`` captures inventory + wave state",
                "+ per-unit status / verdict / tests / score",
            ]),
            ("Resumable", [
                "A second invocation with the same ``REPO_ID`` can resume",
                "from the last wave; already-succeeded units aren't re-run",
            ]),
            ("Workdir snapshots (per-unit)", [
                "Every verify=pass snapshots the unit's ``src/`` tree;",
                "regression rolls back to the latest snapshot",
            ]),
            ("Audit-grade", [
                "Full reconstruction of every agent call, token count,",
                "and cost for compliance + post-mortem",
            ]),
        ],
    )
    add_footer(s, "Part 2  ·  Budget & checkpoints")


def s19_cross_unit_injection():
    s = blank_slide()
    add_title(s, "Cross-unit API injection: the callee's real Java in the caller's prompt",
              subtitle_text="Wave ordering + workdir introspection + explicit FQN directives.  Eliminates hallucinated stubs.  Informed by Anthropic's agentic-loop guidance.",
              section="PART 2 — PLATFORM ARCHITECTURE")
    # Left: how it works. Right: example prompt fragment.
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    top = Inches(1.85)
    h = Inches(5.0)
    card(s, MARGIN, top, col_w, h,
         "How it works",
         [
             ("Step 1  ·  Grapher emits `calls` edges", [
                 "``grapher_hints()`` walks the KG and reports every CALL",
                 "target that resolves to a unit in the inventory",
             ]),
             ("Step 2  ·  Planner orders callees before callers", [
                 "STMTGEN calls INTCALC  →  INTCALC lands in wave 1,",
                 "STMTGEN in wave 2",
             ]),
             ("Step 3  ·  Translator injection", [
                 "For every callee of the current unit, read its",
                 "``units/<CALLEE>/src/main/java/.../service/*.java``",
                 "and ``dto/*.java`` from disk and inline into the prompt",
             ]),
             ("Step 4  ·  Explicit directive in header", [
                 "  ``Inject as:  import com.example.intcalc.service.IntcalcService;``",
                 "  ``This OVERRIDES the generic 'ONE stub per CALL' rule.``",
             ]),
             ("Size-bounded", [
                 "``_MAX_CALLEE_API_BYTES = 40_000``; elides with a marker",
                 "so a giant callee can't blow the prompt context",
             ]),
         ], accent=ACCENT)
    code = """## Already-migrated callee APIs
   (ground truth — use verbatim)

### Callee `INTCALC` — migrated API

Inject as: `import com.example.intcalc.
          service.IntcalcService;`
→ field `private final IntcalcService
          intcalcService;`

#### `src/main/java/com/example/intcalc/
      service/IntcalcService.java`
```java
@Service
public class IntcalcService {
    public IntcalcResult calculateInterest(
        BigDecimal balance, BigDecimal rate,
        BigDecimal accrued, LnkStatus status
    ) { ... }
}
```

#### `.../dto/IntcalcResult.java`
```java
public record IntcalcResult(
    BigDecimal interest, LnkStatus status
) {}
```"""
    code_block(s, MARGIN + col_w + Inches(0.3), top, col_w, h, code, font_size=Pt(10))
    add_footer(s, "Part 2  ·  Cross-unit API injection")


def s20_cartridge_aws_scope():
    s = blank_slide()
    add_title(s, "AWS Lambda (polyglot) → Azure Functions (Python)",
              subtitle_text="Pragmatic consolidation: five Lambda runtimes collapse onto one Azure Functions target.",
              section="PART 3 — SAMPLE CARTRIDGE 1")
    two_column(
        s,
        left_title="Source: AWS Lambda estate",
        left_items=[
            ("Five handler languages", [
                "Python 3.x  ·  Node.js 20  ·  TypeScript",
                "Java 17  ·  C# / .NET",
            ]),
            ("Common shape, divergent idioms", [
                "``def handler(event, context)`` / ``exports.handler`` /",
                "``public class Handler implements RequestHandler<I,O>``",
            ]),
            ("Wired to the AWS ecosystem", [
                "``boto3`` / AWS SDK for JS / V2 / .NET",
                "SAM templates + API Gateway + DynamoDB + S3",
            ]),
            ("Tribal deployments", [
                "SAM YAML and serverless.yml side-by-side;",
                "per-environment hand edits",
            ]),
        ],
        right_title="Target: Azure Functions (Python v2 model)",
        right_items=[
            ("One runtime, one language", [
                "Python 3.11 Azure Functions v2 (decorator model)",
                "``@app.function_name`` / ``@app.route`` / timer / queue",
            ]),
            ("Azure ecosystem equivalents", [
                "``azure-storage-blob``  ←  S3",
                "``azure-cosmos``  ←  DynamoDB",
                "``azure-servicebus``  ←  SNS / SQS",
            ]),
            ("Infra as code", [
                "SAM + serverless.yml mined for env/binding hints;",
                "Bicep is a future step (out of scope for v1)",
            ]),
            ("Correctness signal", [
                "Smoke-boot via Functions Core Tools + unit tests;",
                "dual-run against captured Lambda invocations",
            ]),
        ],
    )
    add_footer(s, "Part 3  ·  Cartridge scope")


def s21_cartridge_aws_recipes():
    s = blank_slide()
    add_title(s, "Recipes do the mechanical work; residuals go to the LLM",
              subtitle_text="Deterministic transforms shrink the prompt surface. The LLM only sees what static tools couldn't fix.",
              section="PART 3 — SAMPLE CARTRIDGE 1")
    two_column(
        s,
        left_title="Deterministic recipes (per language adapter)",
        left_items=[
            ("Python — libcst transformers", [
                "``BotoImportToAzureStub``",
                "  boto3 import → azure-SDK stub + TODO",
                "``HandlerShapeRewrite``",
                "  ``def handler(event, context)`` → ``def main(req:",
                "  func.HttpRequest) -> func.HttpResponse``",
            ]),
            ("Node / TypeScript — tree-sitter", [
                "``aws-sdk`` import rewriter",
                "``exports.handler`` → ``@app.function_name`` binding",
            ]),
            ("Java / C# — annotation rewrites", [
                "Handler class decoration → Azure Functions",
                "attribute-driven runtime",
            ]),
            ("SAM / serverless.yml — IaC ingest", [
                "Parses event bindings (API Gateway, S3, DynamoDB);",
                "surfaces them as translator prompt hints",
            ]),
        ],
        right_title="Residuals: what's left for the LLM",
        right_items=[
            ("AWS imports the recipes couldn't stub", [
                "Third-party wrappers, custom libraries",
            ]),
            ("Handler bodies with inline AWS calls", [
                "`boto3.resource('dynamodb')` mid-function",
                "`s3.upload_file(..., StorageClass='GLACIER')`",
            ]),
            ("Output shape mismatches", [
                "API Gateway proxy response envelope →",
                "``func.HttpResponse`` body + status + headers",
            ]),
            ("What the translator sees", [
                "Residuals markdown (what's unresolved),",
                "plus the file's post-recipe text,",
                "plus IaC binding hints",
            ]),
        ],
    )
    add_footer(s, "Part 3  ·  Recipes & residuals")


def s22_cartridge_aws_example():
    s = blank_slide()
    add_title(s, "Example: Python Lambda → Python Azure Function",
              subtitle_text="Handler rewrite + boto3 stubbing + binding inference, then LLM-driven polish.",
              section="PART 3 — SAMPLE CARTRIDGE 1")
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    before = """# Source: AWS Lambda (Python)
import json, boto3
from decimal import Decimal

dynamo = boto3.resource("dynamodb")
table = dynamo.Table("Customers")

def handler(event, context):
    customer_id = event["pathParameters"]["id"]
    resp = table.get_item(Key={"id": customer_id})
    item = resp.get("Item")
    if not item:
        return {"statusCode": 404,
                "body": "not found"}
    return {"statusCode": 200,
            "body": json.dumps(item,
                               default=str)}"""
    after = """# Target: Azure Functions v2 (Python)
import json
import azure.functions as func
from azure.cosmos import CosmosClient

app = func.FunctionApp()
client = CosmosClient.from_connection_string(
    os.environ["COSMOS_CONN"]
)
container = client.get_database_client("app") \\
                  .get_container_client("Customers")

@app.function_name("get_customer")
@app.route(route="customers/{id}",
           methods=[func.HttpMethod.GET])
def main(req: func.HttpRequest) -> func.HttpResponse:
    customer_id = req.route_params.get("id")
    try:
        item = container.read_item(
            item=customer_id, partition_key=customer_id,
        )
    except CosmosResourceNotFoundError:
        return func.HttpResponse("not found",
                                 status_code=404)
    return func.HttpResponse(
        body=json.dumps(item, default=str),
        status_code=200,
        mimetype="application/json",
    )"""
    code_block(s, MARGIN, Inches(1.85), col_w, Inches(5), before, font_size=Pt(9.5))
    code_block(s, MARGIN + col_w + Inches(0.3), Inches(1.85), col_w, Inches(5), after, font_size=Pt(9.5))
    add_textbox(
        s, MARGIN, Inches(6.95), SLIDE_W - 2 * MARGIN, Inches(0.3),
        [("Recipes handled the import + handler shape.  The LLM drove the DynamoDB→Cosmos mapping, route binding, and error translation.",
          {"size": Pt(11), "italic": True, "color": GRAY,
           "align": PP_ALIGN.CENTER})],
    )
    add_footer(s, "Part 3  ·  Before / after")


def s23_cartridge_aws_review():
    s = blank_slide()
    add_title(s, "Reviewer rules: AWS→Azure-specific findings",
              subtitle_text="A subset of the cartridge's reviewer rule catalog. Each is machine-enforceable + error-level.",
              section="PART 3 — SAMPLE CARTRIDGE 1")
    add_table(
        s, MARGIN, Inches(1.85), SLIDE_W - 2 * MARGIN, Inches(5.0),
        ["Rule slug", "Fires when", "Verdict"],
        [
            ("aws-sdk-import-leaked",       "An `import boto3` / `aws-sdk` / `amazon.*` survives in the emitted code.", "revise"),
            ("lambda-handler-shape",        "A function with signature `(event, context)` is emitted as the entry point.", "revise"),
            ("missing-azure-function-binding", "No `@app.function_name` / binding annotation on the handler.", "revise"),
            ("cosmos-partition-key-missing","`container.read_item` / `delete_item` called without `partition_key`.", "revise"),
            ("apigw-proxy-envelope-leaked", "Return shape still has `statusCode` / `body` keys rather than `func.HttpResponse`.", "revise"),
            ("connection-string-inline",    "Azure connection string / key embedded in source instead of `os.environ[...]`.", "revise"),
            ("dynamodb-decimal-residual",   "`decimal.Decimal` handling from the Lambda kept but Cosmos returns native JSON.", "warning"),
        ],
        body_font=Pt(10),
    )
    add_footer(s, "Part 3  ·  Reviewer rules")


def s30_cobol_scope():
    s = blank_slide()
    add_title(s, "COBOL + JCL + DB2 + CICS → Java 25 LTS / Spring Boot 4.0.5",
              subtitle_text="The hard one: shape-shifting source, multiple target styles, and correctness you can prove.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    two_column(
        s,
        left_title="Source: mainframe COBOL estate",
        left_items=[
            ("COBOL language features we handle", [
                "Fixed-format + free-format",
                "REDEFINES, 88-level, OCCURS DEPENDING ON",
                "COMP-3 / PACKED-DECIMAL, PIC with S, V",
                "Copybook (`.cpy`) expansion + grouping",
            ]),
            ("Batch & transactional contexts", [
                "JCL: jobs, steps, DD cards → datasets",
                "EXEC SQL (DB2) blocks — CRUD classification",
                "EXEC CICS — COMMAREA, maps, transaction IDs",
            ]),
            ("Inter-program coupling", [
                "`CALL 'TARGET' USING …` — cross-unit CALL graph",
                "LINKAGE SECTION shapes the public method signature",
            ]),
        ],
        right_title="Target: Java 25 LTS + Spring Boot 4.0.5",
        right_items=[
            ("Three target styles, one per program", [
                "batch       → Spring Batch (Job + Tasklets + Readers + Writers)",
                "cics        → Spring MVC (REST endpoint per tx)",
                "subroutine  → @Service (LINKAGE → method signature)",
            ]),
            ("Modern Java idioms, mandatory", [
                "``record`` DTOs, sealed interfaces, pattern switch,",
                "text blocks for SQL, BigDecimal + MathContext",
            ]),
            ("Compile-gated", [
                "``mvn compile`` must succeed before the unit is accepted",
                "Parallel dual-run gate for equivalence",
            ]),
            ("Toolchain pinned", [
                "`<java.version>25</java.version>` +",
                "`spring-boot-starter-parent 4.0.5`  (reviewer-enforced)",
            ]),
        ],
    )
    add_footer(s, "Part 4  ·  Scope")


def s31_cobol_adapter():
    s = blank_slide()
    add_title(s, "COBOL adapter: structure-preserving parse, then KG population",
              subtitle_text="tree-sitter + regex fallback. REDEFINES / 88-levels / OCCURS are first-class KG citizens.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    code = """# adapters/cobol.py — what lands in the KG per field
# ------------------------------------------------------------
#   01  CUSTOMER-REC.
#       05  CUST-ID        PIC 9(10).
#       05  CUST-BAL       PIC S9(13)V9(2) COMP-3.
#       05  CUST-STATUS    PIC X(1).
#           88  ACTIVE     VALUE 'A'.
#           88  CLOSED     VALUE 'C'.
#       05  CUST-SSN-RAW   PIC 9(9).
#       05  CUST-SSN-FMT   REDEFINES CUST-SSN-RAW PIC X(11).
#
# KG nodes:
#   field:CUSTOMER-REC          kind=field (group, parent)
#   field:CUST-ID               kind=field  pic=9(10)         numeric
#   field:CUST-BAL              kind=field  pic=S9(13)V9(2)   packed
#                               scale=2  mathcontext=required
#   field:CUST-STATUS           kind=field  pic=X(1)
#     (conditional) ACTIVE  = 'A'
#     (conditional) CLOSED  = 'C'
#   field:CUST-SSN-RAW          kind=field  pic=9(9)
#   field:CUST-SSN-FMT          kind=field  pic=X(11)
#     --redefines--> field:CUST-SSN-RAW
#
# Callable subroutines:
#   `PROCEDURE DIVISION USING LNK-BALANCE LNK-RATE ...`
#       → field:LNK-BALANCE  --input_of--> program:INTCALC
#       → field:LNK-RATE     --input_of--> program:INTCALC"""
    code_block(s, MARGIN, Inches(1.85), SLIDE_W - 2 * MARGIN, Inches(5.0), code, font_size=Pt(11))
    add_footer(s, "Part 4  ·  COBOL adapter")


def s32_cobol_target_style():
    s = blank_slide()
    add_title(s, "Target-style classifier: one prompt per shape",
              subtitle_text="Signals from the KG decide whether a program becomes a Batch Job, a CICS REST controller, or a Service bean.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    add_table(
        s, MARGIN, Inches(1.9), SLIDE_W - 2 * MARGIN, Inches(3.3),
        ["Target style", "Signals that pick it", "Generated shape", "Translator prompt"],
        [
            ("cics",       "≥1 `txn_call` descendants (EXEC CICS). Wins over batch if both apply.",
             "Spring MVC `@RestController` + `@Service` + JPA repo",
             "translator_cics.md"),
            ("subroutine", "`PROCEDURE DIVISION USING …` present (LINKAGE). No JCL wrapper, no CICS.",
             "`@Service` bean; LINKAGE → public method signature",
             "translator_subroutine.md"),
            ("batch",      "`step_of` edge from a JCL step to the program. No LINKAGE, no CICS.",
             "Spring Batch `@Configuration` with Job + Tasklets + Readers + Writers",
             "translator_batch.md"),
        ],
        body_font=Pt(11),
    )
    card(s, MARGIN, Inches(5.35), SLIDE_W - 2 * MARGIN, Inches(1.5),
         "Why the split matters",
         [
             ("Different deliverables per style", [
                 "Batch must emit no `@RestController`; CICS must emit no `Job` / `Tasklet`;",
                 "subroutine must emit no `Main` that runs on startup. The translator prompt",
                 "enforces these as DO NOT rules; the reviewer has a `wrong-target-style` rule.",
             ]),
         ], accent=PRIMARY)
    add_footer(s, "Part 4  ·  Target-style classifier")


def s33_cobol_kg_example():
    s = blank_slide()
    add_title(s, "KG example: the pilot corpus (CUSTINQ / INTCALC / STMTGEN)",
              subtitle_text="Three programs, three target styles, one CALL edge that shapes wave ordering.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    # Three program boxes + dataset pills + edges implied via text
    top = Inches(1.85)
    box_w = Inches(3.6)
    box_h = Inches(1.8)
    gap = Inches(0.4)
    total = 3 * box_w + 2 * gap
    start_x = (SLIDE_W - total) / 2
    programs = [
        ("CUSTINQ",   "cics",       "EXEC CICS RECEIVE / SEND\nEXEC SQL SELECT  ←  Customer",   ACCENT),
        ("INTCALC",   "subroutine", "LINKAGE SECTION\nBigDecimal interest calc",                 GREEN),
        ("STMTGEN",   "batch",      "JCL step STEP010\nReads CUST-MASTER  →  STMT-OUT\nCALL 'INTCALC'  on each customer", PRIMARY),
    ]
    for i, (pid, style, body, color) in enumerate(programs):
        x = start_x + (box_w + gap) * i
        add_rect(s, x, top, box_w, box_h,
                 fill=LIGHT, line=color, line_weight=Pt(2),
                 text=[(pid, {"size": Pt(18), "bold": True, "color": color,
                              "align": PP_ALIGN.CENTER, "mono": True}),
                       (f"target style: {style}",
                        {"size": Pt(11), "color": GRAY,
                         "align": PP_ALIGN.CENTER, "italic": True}),
                       ("", {"size": Pt(4)}),
                       (body, {"size": Pt(10), "color": DARK,
                               "align": PP_ALIGN.CENTER})],
                 text_kwargs={"anchor": MSO_ANCHOR.MIDDLE})
    # Arrow STMTGEN  → INTCALC
    arrow_y = top + box_h + Inches(0.3)
    add_textbox(
        s, MARGIN, arrow_y, SLIDE_W - 2 * MARGIN, Inches(0.5),
        [("STMTGEN  —calls→  INTCALC      (grapher_hints emits this; planner schedules INTCALC in wave 1, STMTGEN in wave 2)",
          {"size": Pt(13), "color": PRIMARY, "bold": True,
           "align": PP_ALIGN.CENTER, "mono": True})],
    )
    # Bottom: the CRUD matrix snippet
    card(s, MARGIN, Inches(5.0), SLIDE_W - 2 * MARGIN, Inches(1.9),
         "CRUD matrix (program × data asset)",
         [
             ("Customer table (DB2)", [
                 "CUSTINQ = R       STMTGEN = R       INTCALC = —",
             ]),
             ("CUST-MASTER  /  STMT-OUT (flat files via JCL DD)", [
                 "STMTGEN = R       STMTGEN = W       CUSTINQ = —",
             ]),
             ("Derived seam ranking (Thoughtworks methodology)", [
                 "CUSTINQ + Customer  →  strong-cohesion candidate for REST API slice",
                 "STMTGEN + CUST-MASTER / STMT-OUT  →  batch pipeline, isolated",
             ]),
         ], accent=GREEN)
    add_footer(s, "Part 4  ·  KG example")


def s34_cobol_business_spec():
    s = blank_slide()
    add_title(s, "Business spec: what the LLM actually sees",
              subtitle_text="Deterministic render from the KG + summaries. Per-unit, target-style aware.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    spec = """# Business spec — STMTGEN  (target style: batch)

## Purpose
Generate the monthly customer statement file. For every record in
CUST-MASTER, compute accrued interest via CALL 'INTCALC' and
write one STMT-OUT record per customer.

## Inputs
- File CUST-MASTER  (FD CUST-MASTER; DD CUSTIN on PROD.CUST.MASTER)
- DB2 table CUSTOMER (EXEC SQL SELECT balance, rate, accrued INTO ...)

## Outputs
- File STMT-OUT     (FD STMT-OUT; DD STMTOUT on PROD.STMT.DAILY)
- DB2 table STATEMENT_HISTORY (INSERT)

## Side effects
- CALL 'INTCALC' USING LNK-BALANCE LNK-RATE LNK-INTEREST LNK-STATUS

## Key invariants
- Every input customer produces exactly one output STMT-OUT record
- Interest calculation uses MathContext(18, HALF_UP) scale=2
- ABEND on negative balance: LNK-STATUS='N' propagates to exit code

## Data dictionary
| Field         | PIC           | Java type   | Notes                          |
| CUST-ID       | 9(10)         | Long        | primary key                    |
| CUST-BAL      | S9(13)V9(2)   | BigDecimal  | packed; scale=2                |
| CUST-RATE     | S9(1)V9(9)    | BigDecimal  | packed; scale=9, range [0,1]   |
| CUST-STATUS   | X(1)          | enum        | conditionals ACTIVE/CLOSED     |"""
    code_block(s, MARGIN, Inches(1.85), SLIDE_W - 2 * MARGIN, Inches(5.0), spec, font_size=Pt(10.5))
    add_footer(s, "Part 4  ·  Business spec rendering")


def s35_cobol_translator_prompt():
    s = blank_slide()
    add_title(s, "Translator prompt: structure, not verbosity",
              subtitle_text="One prompt per target style. Pre-computed Java naming removes an entire class of LLM errors.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    two_column(
        s,
        left_title="What every translator prompt includes",
        left_items=[
            ("Pre-computed Java naming block", [
                "Exact package, class_base, file paths per deliverable",
                "Translator can't mis-capitalise ``CustinqService`` again",
            ]),
            ("Style-specific deliverables list", [
                "What MUST be emitted + what MUST NOT (DO NOT rules)",
            ]),
            ("pom.xml contract — toolchain pinned", [
                "`<java.version>25</java.version>` only",
                "`spring-boot-starter-parent 4.0.5`  (Spring Framework 7)",
            ]),
            ("DTO-as-record mandate", [
                "`public record FooDto(...) {}` — eliminates twin-DTO class",
            ]),
            ("Cross-file consistency rules", [
                "No twin-record DTO, no sealed-without-implements,",
                "no duplicate external stub, no `javax.*` imports",
            ]),
        ],
        right_title="Cross-unit injection block (when callees exist)",
        right_items=[
            ("Appended after the base prompt", [
                "Read callee's already-emitted service + DTOs from disk",
                "Inline verbatim with explicit `import <FQN>` directive",
            ]),
            ("Overrides the generic ``ONE stub per CALL`` rule", [
                "For callees in the injection block → ZERO stubs",
                "Re-declaring them is the `duplicate-external-stub` finding",
            ]),
            ("Reviewer backs this up", [
                "`cross-module-reimplementation` rule (in roadmap) greps",
                "for the callee's class names in sub-packages of the caller",
            ]),
            ("Size-bounded", [
                "40 KB ceiling with elision marker when exceeded",
            ]),
        ],
    )
    add_footer(s, "Part 4  ·  Translator prompt strategy")


def s36_cobol_reviewer_rules():
    s = blank_slide()
    add_title(s, "Reviewer catalog (abbreviated)",
              subtitle_text="Every rule has a slug, severity, and verdict. The retry loop merges findings into the next translator prompt.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    add_table(
        s, MARGIN, Inches(1.85), SLIDE_W - 2 * MARGIN, Inches(5.0),
        ["Slug", "Category", "What it catches", "Verdict"],
        [
            ("bigdecimal-no-mathcontext",   "numeric",     "BigDecimal arithmetic without explicit MathContext", "revise"),
            ("primitive-float-for-numeric", "numeric",     "double/float used for fields sourced from COMP-3 / PIC 9…V9…", "revise"),
            ("linkage-signature-mismatch",  "structural",  "Subroutine's public method doesn't match LINKAGE order/types", "revise"),
            ("missing-required-deliverable","structural",  "pom.xml / Application / ContextLoadTest / yml / README absent", "revise"),
            ("undefined-class-reference",   "cross-file",  "A type used in one file is never defined in any emitted file", "revise"),
            ("undefined-method-reference",  "cross-file",  "A method call on an injected service that isn't declared", "revise"),
            ("twin-record-dto",             "cross-file",  "Both `Foo` and `FooRecord` emitted in the same module", "revise"),
            ("sealed-permits-without-implements", "cross-file", "`sealed interface X permits A, B` with A/B not implementing X", "revise"),
            ("duplicate-external-stub",     "cross-file",  "Multiple service-shaped classes for one COBOL CALL target", "revise"),
            ("spring-batch-api-mismatch",   "api",         "ExecutionContext.putBoolean(...), wrong Tasklet return shape, etc.", "revise"),
            ("java-version-below-25",       "toolchain",   "pom.xml declares `<java.version>` below 25", "revise"),
            ("spring-boot-parent-below-4.0","toolchain",   "`spring-boot-starter-parent` version is 3.x", "revise"),
            ("javax-import-on-jakarta-runtime", "toolchain", "Any `import javax.*` (should be `jakarta.*` under Spring Boot 4)", "revise"),
            ("sql-concatenated",            "security",    "String-concatenated SQL instead of named `@Query` parameters", "revise"),
            ("field-injection",             "security",    "`@Autowired` field or bare field assignment — must be constructor injection", "revise"),
            ("hardcoded-config",            "security",    "Connection strings / secrets inline instead of `application.yml` / env", "revise"),
        ],
        body_font=Pt(9.5), header_font=Pt(11),
    )
    add_footer(s, "Part 4  ·  Reviewer catalog")


def s37_cobol_results():
    s = blank_slide()
    add_title(s, "Pilot run on the CUSTINQ / INTCALC / STMTGEN corpus",
              subtitle_text="Latest end-to-end run on the three-program pilot with OpenAI gpt-5.4 and cross-unit API injection enabled.",
              section="PART 4 — SAMPLE CARTRIDGE 2")
    # Top row: two big metrics cards
    m_top = Inches(1.85)
    m_h = Inches(1.4)
    m_w = Inches(3.0)
    gap = Inches(0.25)
    start_x = (SLIDE_W - 4 * m_w - 3 * gap) / 2
    metrics = [
        ("3 / 3", "programs succeeded", GREEN),
        ("4m 57s", "total wall time", ACCENT),
        ("$0.4980", "LLM cost / run", PRIMARY),
        ("288", "tests passing", GREEN),
    ]
    for i, (val, label, col) in enumerate(metrics):
        x = start_x + (m_w + gap) * i
        add_rect(s, x, m_top, m_w, m_h,
                 fill=LIGHT, line=col, line_weight=Pt(2),
                 text=[(val, {"size": Pt(28), "bold": True, "color": col,
                              "align": PP_ALIGN.CENTER}),
                       (label, {"size": Pt(11), "color": GRAY,
                                "align": PP_ALIGN.CENTER})],
                 text_kwargs={"anchor": MSO_ANCHOR.MIDDLE})
    # Detail table
    add_table(
        s, MARGIN, Inches(3.55), SLIDE_W - 2 * MARGIN, Inches(3.2),
        ["Program", "Target style", "Java files emitted", "mvn compile", "Reviewer verdict", "Cost"],
        [
            ("CUSTINQ", "cics",       "12", "BUILD SUCCESS", "accept", "$0.0695"),
            ("INTCALC", "subroutine",  "4", "BUILD SUCCESS", "revise", "$0.1180"),
            ("STMTGEN", "batch",      "22", "BUILD SUCCESS", "revise", "$0.3106"),
        ],
        body_font=Pt(12), header_font=Pt(12),
    )
    add_footer(s, "Part 4  ·  Pilot results")


def s40_validation():
    s = blank_slide()
    add_title(s, "Validation: compile gate + dual-run equivalence",
              subtitle_text="Static type-checking catches hallucinations. Behavioural equivalence catches everything else.",
              section="PART 5 — VALIDATION & ROADMAP")
    two_column(
        s,
        left_title="Compile gate (fast, cheap)",
        left_items=[
            ("`mvn compile` per unit", [
                "Runs after every translator attempt + every retry",
                "stderr tail fed back into the next retry as feedback",
            ]),
            ("Verify feedback merged into prompt", [
                "`verify_feedback.md` lives in the unit's workdir;",
                "reviewer findings + compiler errors combine",
            ]),
            ("Halt if never green", [
                "After `max_reviewer_retries` failed attempts the unit is",
                "marked `failed` and the next wave skips its dependents",
            ]),
            ("In production runs", [
                "3 / 3 units compile clean on pilot corpus;",
                "`BUILD SUCCESS` is the minimum acceptance bar",
            ]),
        ],
        right_title="Dual-run harness (when fixtures exist)",
        right_items=[
            ("GnuCOBOL  ×  JVM", [
                "Record golden fixtures by replaying prod inputs through",
                "a GnuCOBOL-compiled binary; canonicalise outputs",
            ]),
            ("Tolerance model", [
                "COBOL-aware: rounding, packed-decimal drift,",
                "timestamp normalisation, SQL row-order tolerance",
            ]),
            ("Verdict states", [
                "PASS / MINOR_DRIFT / FAIL; MINOR_DRIFT surfaces diffs",
                "without blocking the unit",
            ]),
            ("Harness is pluggable", [
                "`Runner` ABC; any cartridge can plug a dual-run runner",
                "(`SubprocessRunner` is the reference impl)",
            ]),
        ],
    )
    add_footer(s, "Part 5  ·  Validation")


def s41_roadmap():
    s = blank_slide()
    add_title(s, "Roadmap: where the platform is heading",
              subtitle_text="Near-term hardening + the next two cartridges.",
              section="PART 5 — VALIDATION & ROADMAP")
    col_w = (SLIDE_W - 2 * MARGIN - 2 * Inches(0.3)) / 3
    top = Inches(1.85)
    h = Inches(4.9)
    card(s, MARGIN, top, col_w, h,
         "Near-term (platform)",
         [
             ("Per-request LLM timeout", [
                 "Prevent the 24-min Anthropic retry-storm hang we saw",
                 "in one pilot run",
             ]),
             ("`cross-module-reimplementation` reviewer rule", [
                 "Close the sub-package loophole the LLM used to evade",
                 "`duplicate-external-stub`",
             ]),
             ("GnuCOBOL dual-run wiring", [
                 "Record fixtures + enable default dual-run gate",
                 "for the COBOL cartridge",
             ]),
             ("CI on every cartridge", [
                 "Matrix build per cartridge against a representative corpus",
             ]),
         ], accent=ACCENT)
    card(s, MARGIN + col_w + Inches(0.3), top, col_w, h,
         "Near-term (COBOL cartridge)",
         [
             ("Multi-module Maven support", [
                 "Today each unit is a self-contained module.",
                 "Produce a parent POM so callees are real classpath deps",
             ]),
             ("Scale-up from 3 → AWS Card Demo (~40 programs)", [
                 "Validate wave parallelism at real size",
             ]),
             ("CICS COMMAREA / BMS map parser", [
                 "Richer input for CICS target style",
             ]),
             ("Copybook scope analysis", [
                 "Detect + refactor copybook fan-out patterns",
             ]),
         ], accent=PRIMARY)
    card(s, MARGIN + 2 * col_w + 2 * Inches(0.3), top, col_w, h,
         "Next cartridges",
         [
             ("Legacy .NET Framework → .NET 9", [
                 "WCF + Remoting → modern WebApi + gRPC",
             ]),
             ("PL/SQL → Postgres + pg_plsql", [
                 "Oracle stored-procs as first-class migration target",
             ]),
             ("VB6 → .NET / React", [
                 "Form-shape lift with explicit validation carry-over",
             ]),
             ("Spring Boot 2.x → Spring Boot 4.0.5", [
                 "In-ecosystem uplift; Jakarta rename + deprecated API fix",
             ]),
         ], accent=GREEN)
    add_footer(s, "Part 5  ·  Roadmap")


def s41b_references():
    s = blank_slide()
    add_title(s, "References",
              subtitle_text="Sources of record for every influence called out in this deck.",
              section="PART 5 — VALIDATION & ROADMAP")
    # One card per reference, stacked vertically for readability
    refs = [
        (
            "Martin Fowler — Legacy Modernization meets Generative AI",
            "martinfowler.com, Lilly Ryan, Shodhan Sheth, et al.",
            "https://martinfowler.com/articles/legacy-modernization-gen-ai.html",
            "Our direct blueprint for the KG + post-order DFS summariser + "
            "business-spec rendering.  The CRUD matrix and seam-ranking approach "
            "carry through verbatim.",
        ),
        (
            "Thoughtworks — Claude Code for COBOL Modernization: a reality check",
            "thoughtworks.com/insights (2025)",
            "https://www.thoughtworks.com/en-gb/insights/articles/"
            "claude-code-cobol-modernization-reality",
            "Grounds the design in real-world agentic runs on mainframe corpora.  "
            "Shaped our target-style split (batch / CICS / subroutine), our insistence "
            "on mvn compile as the acceptance gate, and the structured-spec-over-raw-code "
            "prompt discipline.",
        ),
        (
            "Anthropic — How AI helps break the cost barrier for COBOL modernization",
            "claude.com/blog",
            "https://claude.com/blog/how-ai-helps-break-cost-barrier-cobol-modernization",
            "The cost-per-program narrative and the agentic-retry-loop pattern.  "
            "Our BudgetTracker + MAX_BUDGET_USD + per-model accounting trace directly "
            "to this reference.",
        ),
        (
            "Fowler canon — Refactoring  ·  Feathers — Working Effectively with Legacy Code",
            "martinfowler.com  ·  O'Reilly",
            "https://martinfowler.com/books/refactoring.html",
            "Seams, characterisation tests, avoid-premature-abstraction.  "
            "Foundations behind the cartridge isolation model and the dual-run harness.",
        ),
    ]
    top = Inches(1.85)
    card_h = Inches(1.18)
    gap = Inches(0.08)
    for title, byline, url, takeaway in refs:
        add_rect(s, MARGIN, top, SLIDE_W - 2 * MARGIN, card_h,
                 fill=LIGHT, line=BORDER)
        # Title strip
        add_rect(s, MARGIN, top, Inches(0.12), card_h,
                 fill=ACCENT, line=ACCENT)
        add_textbox(
            s, MARGIN + Inches(0.3), top + Inches(0.1),
            SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(0.35),
            [(title, {"size": Pt(13), "bold": True, "color": PRIMARY})],
        )
        add_textbox(
            s, MARGIN + Inches(0.3), top + Inches(0.4),
            SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(0.3),
            [(f"{byline}   ·   {url}",
              {"size": Pt(10), "color": ACCENT, "mono": True})],
        )
        add_textbox(
            s, MARGIN + Inches(0.3), top + Inches(0.7),
            SLIDE_W - 2 * MARGIN - Inches(0.4), Inches(0.45),
            [(takeaway, {"size": Pt(10.5), "color": DARK})],
        )
        top += card_h + gap
    add_footer(s, "Part 5  ·  References")


def s42_close():
    s = blank_slide()
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(7.5), fill=PRIMARY, line=PRIMARY)
    add_rect(s, Inches(0), Inches(0.4), SLIDE_W, Inches(0.08), fill=ACCENT, line=ACCENT)
    add_textbox(
        s, MARGIN, Inches(1.5), Inches(12), Inches(0.5),
        [("THE PLATFORM IN ONE LINE",
          {"size": Pt(14), "color": RGBColor(0xBB, 0xD0, 0xE6), "bold": True})],
    )
    add_textbox(
        s, MARGIN, Inches(2.1), Inches(12), Inches(2.5),
        [("Parse deterministically,",
          {"size": Pt(34), "color": RGBColor(0xFF, 0xFF, 0xFF), "bold": True}),
         ("translate with LLMs,",
          {"size": Pt(34), "color": RGBColor(0xFF, 0xFF, 0xFF), "bold": True}),
         ("review by rules,",
          {"size": Pt(34), "color": RGBColor(0xFF, 0xFF, 0xFF), "bold": True}),
         ("validate by compilation and dual-run.",
          {"size": Pt(34), "color": RGBColor(0xFF, 0xFF, 0xFF), "bold": True})],
    )
    add_textbox(
        s, MARGIN, Inches(5.8), Inches(12), Inches(0.5),
        [("Two cartridges in, many more to go.",
          {"size": Pt(18), "color": RGBColor(0xBB, 0xD0, 0xE6),
           "italic": True})],
    )
    add_textbox(
        s, MARGIN, Inches(7.05), Inches(12), Inches(0.3),
        [("Legacy Modernization Platform  ·  Design Deck",
          {"size": Pt(10), "color": RGBColor(0xBB, 0xD0, 0xE6)})],
    )


# ---------------------------------------------------------------------- #
# Assemble
# ---------------------------------------------------------------------- #


def build():
    s00_cover()
    s01_contents()
    s02_problem()
    s03_thesis()
    s04_influences()
    s10_architecture()
    s11_cartridge_model()
    s12_discovery()
    s13_knowledge_graph()
    s14_comprehension()
    s15_planning()
    s16_worker()
    s17_agents()
    s18_budget_checkpoint()
    s19_cross_unit_injection()
    s20_cartridge_aws_scope()
    s21_cartridge_aws_recipes()
    s22_cartridge_aws_example()
    s23_cartridge_aws_review()
    s30_cobol_scope()
    s31_cobol_adapter()
    s32_cobol_target_style()
    s33_cobol_kg_example()
    s34_cobol_business_spec()
    s35_cobol_translator_prompt()
    s36_cobol_reviewer_rules()
    s37_cobol_results()
    s40_validation()
    s41_roadmap()
    s41b_references()
    s42_close()

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
